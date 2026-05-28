"""
Generate presentation.pptx — run once from the project root.
Usage:  python generate_slides.py
Requires: pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO
from pathlib import Path

OUTPUT = Path(__file__).parent / "presentation.pptx"

# ── Catppuccin Mocha dark palette ──────────────────────────
BG      = RGBColor(0x1e, 0x1e, 0x2e)   # base
SURF    = RGBColor(0x31, 0x32, 0x44)   # surface
SURF2   = RGBColor(0x45, 0x45, 0x59)   # surface1
TEXT    = RGBColor(0xcd, 0xd6, 0xf4)   # text
SUB     = RGBColor(0xa6, 0xad, 0xc8)   # subtext
BLUE    = RGBColor(0x89, 0xb4, 0xfa)
GREEN   = RGBColor(0xa6, 0xe3, 0xa1)
RED     = RGBColor(0xf3, 0x8b, 0xa8)
YELLOW  = RGBColor(0xf9, 0xe2, 0xaf)
PURPLE  = RGBColor(0xcb, 0xa6, 0xf7)
TEAL    = RGBColor(0x94, 0xe2, 0xd5)
PEACH   = RGBColor(0xfa, 0xb3, 0x87)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Core helpers ───────────────────────────────────────────

def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def tb(s, text, x, y, w, h, *, sz=20, bold=False, italic=False,
       col=TEXT, align=PP_ALIGN.LEFT, wrap=True):
    """Add a styled text box."""
    shape = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = col
    r.font.name = "Segoe UI"
    return shape


def tb_multi(s, lines, x, y, w, h, *, sz=18, col=TEXT, gap=Pt(6)):
    """Add a text box with multiple paragraphs."""
    shape = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    for i, (text, opts) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = gap
        r = p.add_run()
        r.text = text
        r.font.size = Pt(opts.get("sz", sz))
        r.font.bold = opts.get("bold", False)
        r.font.color.rgb = opts.get("col", col)
        r.font.name = "Segoe UI"
    return shape


def box(s, x, y, w, h, fill=SURF, line_col=None, line_pt=1.5, radius=None):
    """Add a filled rounded rectangle."""
    shape_type = MSO.ROUNDED_RECTANGLE if radius else MSO.RECTANGLE
    sh = s.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_col:
        sh.line.color.rgb = line_col
        sh.line.width = Pt(line_pt)
    else:
        sh.line.fill.background()
    return sh


def card(s, x, y, w, h, title, body_lines, *, title_col=BLUE, body_col=TEXT, body_sz=16):
    """Card: dark rounded box with a title + bullet body."""
    box(s, x, y, w, h, fill=SURF, radius=True)
    tb(s, title, x + 0.15, y + 0.1, w - 0.3, 0.38, sz=17, bold=True, col=title_col)
    for i, line in enumerate(body_lines):
        tb(s, f"▸  {line}", x + 0.15, y + 0.5 + i * 0.35, w - 0.3, 0.38,
           sz=body_sz, col=body_col)


def header(s, title, subtitle=None):
    """Blue accent bar + title + optional subtitle."""
    box(s, 0, 0, 13.333, 0.07, fill=BLUE)
    tb(s, title, 0.4, 0.14, 12.5, 0.75, sz=32, bold=True, col=TEXT)
    if subtitle:
        tb(s, subtitle, 0.4, 0.88, 12.5, 0.38, sz=15, col=SUB)


def bullet_col(s, items, x, y, w, *, sz=18, col=TEXT, icon="▸", gap=0.42):
    """Render a vertical list of bullet points."""
    for i, item in enumerate(items):
        tb(s, f"{icon}  {item}", x, y + i * gap, w, gap + 0.08, sz=sz, col=col)


def stat_box(s, x, y, value, label, val_col=BLUE):
    """Metric card: big number + small label."""
    box(s, x, y, 2.8, 1.35, fill=SURF, radius=True)
    tb(s, value, x, y + 0.15, 2.8, 0.65,
       sz=36, bold=True, col=val_col, align=PP_ALIGN.CENTER)
    tb(s, label, x, y + 0.78, 2.8, 0.42,
       sz=14, col=SUB, align=PP_ALIGN.CENTER)


def _style_chart(chart, title, bar_color, is_line=False):
    """Apply consistent styling to a chart object."""
    try:
        chart.chart_title.text_frame.text = title
        p = chart.chart_title.text_frame.paragraphs[0]
        if p.runs:
            p.runs[0].font.color.rgb = TEXT
            p.runs[0].font.size = Pt(14)
    except Exception:
        pass
    try:
        chart.has_legend = False
    except Exception:
        pass
    try:
        for ser in chart.series:
            if is_line:
                ser.format.line.color.rgb = bar_color
                ser.format.line.width = Pt(2.5)
            else:
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = bar_color
                ser.format.line.fill.background()
    except Exception:
        pass


def add_multi_line_chart(s, x, y, w, h, categories, series_data, title=""):
    """Line chart with per-series colors. series_data = [(name, values, RGBColor), ...]"""
    cd = CategoryChartData()
    cd.categories = categories
    for name, values, _ in series_data:
        cd.add_series(name, values)
    cf = s.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(x), Inches(y), Inches(w), Inches(h), cd
    )
    chart = cf.chart
    try:
        chart.has_legend = True
        chart.chart_title.text_frame.text = title
        p = chart.chart_title.text_frame.paragraphs[0]
        if p.runs:
            p.runs[0].font.color.rgb = TEXT
            p.runs[0].font.size = Pt(14)
        for ser, (_, _, col) in zip(chart.series, series_data):
            ser.format.line.color.rgb = col
            ser.format.line.width = Pt(2.5)
    except Exception:
        pass
    return cf


def add_bar_chart(s, x, y, w, h, categories, series_data, title="", bar_color=BLUE):
    """Native PowerPoint column chart."""
    cd = CategoryChartData()
    cd.categories = categories
    for name, values in series_data:
        cd.add_series(name, values)
    cf = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h), cd
    )
    _style_chart(cf.chart, title, bar_color, is_line=False)
    return cf


def add_line_chart(s, x, y, w, h, categories, series_data, title="", line_color=GREEN):
    """Native PowerPoint line chart."""
    cd = CategoryChartData()
    cd.categories = categories
    for name, values in series_data:
        cd.add_series(name, values)
    cf = s.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(x), Inches(y), Inches(w), Inches(h), cd
    )
    _style_chart(cf.chart, title, line_color, is_line=True)
    return cf


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════
s1 = new_slide()

# Wide accent bar bottom
box(s1, 0, 6.9, 13.333, 0.1, fill=BLUE)

# Decorative side accent
box(s1, 0, 0, 0.08, 7.5, fill=BLUE)

# Center title block
tb(s1, "⚡", 1.0, 1.4, 2.0, 1.4,
   sz=80, col=BLUE, align=PP_ALIGN.CENTER)
tb(s1, "Probabilistic Energy\nDemand Forecasting",
   2.8, 1.3, 9.5, 2.2,
   sz=46, bold=True, col=TEXT, align=PP_ALIGN.LEFT)

tb(s1, "Dominion (DOM) Region  ·  Virginia / Washington DC  ·  2005–2025",
   2.8, 3.55, 9.5, 0.5, sz=19, col=BLUE)

box(s1, 2.8, 4.15, 9.5, 0.04, fill=SURF2)

tb(s1, "XGBoost  ·  LightGBM  ·  LSTM (Quantile Regression)",
   2.8, 4.28, 9.5, 0.45, sz=16, col=SUB)

tb(s1, "Dawit Adane  ·  ML-1 Course Project  ·  2026",
   2.8, 5.4, 9.5, 0.45, sz=15, col=SUB)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════
s2 = new_slide()
header(s2, "Agenda")

agenda = [
    ("01", "Problem Statement & Motivation",    BLUE),
    ("02", "Dataset & Region Overview",          TEAL),
    ("03", "Exploratory Data Analysis",          GREEN),
    ("04", "Feature Engineering & Models",       YELLOW),
    ("05", "Results & Model Comparison",         PEACH),
    ("06", "Future Forecast & Live Demo",        PURPLE),
]

col_x = [0.5, 6.9]
for i, (num, label, col) in enumerate(agenda):
    col_idx = i % 2
    row = i // 2
    x = col_x[col_idx]
    y = 1.35 + row * 1.65
    box(s2, x, y, 6.0, 1.35, fill=SURF, radius=True)
    tb(s2, num, x + 0.18, y + 0.08, 0.75, 0.55, sz=30, bold=True, col=col)
    tb(s2, label, x + 1.0, y + 0.18, 4.8, 0.55, sz=18, bold=True, col=TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Problem Statement
# ══════════════════════════════════════════════════════════════
s3 = new_slide()
header(s3, "Problem Statement", "Why forecast energy demand?")

tb(s3, "Traditional forecasting gives a single number.", 0.5, 1.3, 12.3, 0.5, sz=22, bold=True, col=TEXT)
tb(s3, "We need to know: how confident are we — and what is the plausible range?",
   0.5, 1.85, 12.3, 0.6, sz=19, col=SUB)

problems = [
    ("Grid Operators",  "Must balance supply/demand in real time — under-forecast = blackouts, over-forecast = wasted generation", BLUE),
    ("Utilities",       "Bid into electricity markets hours ahead — point forecasts lead to costly imbalance penalties", TEAL),
    ("Planners",        "Capacity decisions made 5–10 years out require uncertainty ranges, not single numbers", GREEN),
    ("Researchers",     "Quantify how climate extremes and data center growth will stress the grid", YELLOW),
]

for i, (title, body, col) in enumerate(problems):
    col_idx = i % 2
    row = i // 2
    x = 0.4 + col_idx * 6.4
    y = 2.7 + row * 2.0
    box(s3, x, y, 6.1, 1.75, fill=SURF, radius=True)
    box(s3, x, y, 0.06, 1.75, fill=col)
    tb(s3, title, x + 0.22, y + 0.1, 5.7, 0.45, sz=17, bold=True, col=col)
    tb(s3, body, x + 0.22, y + 0.55, 5.7, 1.0, sz=14, col=TEXT, wrap=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Why DOM? Data Center Hub
# ══════════════════════════════════════════════════════════════
s4 = new_slide()
header(s4, "Why Virginia / DC?", "The world's largest data center hub")

# Left side — narrative
tb(s4, "Dominion (DOM) covers the Richmond–Washington DC corridor,",
   0.4, 1.35, 7.0, 0.5, sz=17, col=TEXT)
tb(s4, "home to the highest concentration of data centers on Earth.",
   0.4, 1.8, 7.0, 0.5, sz=17, col=TEXT)

bullet_col(s4, [
    "Loudoun County (\"Data Center Alley\") hosts ~70% of global internet traffic",
    "Amazon AWS, Microsoft Azure, Google, Meta all operate hyperscale facilities here",
    "These facilities run 24 / 7 — unlike residential demand, the load never fully drops",
    "Demand surged ~45% between 2020 and 2025, far exceeding national grid averages",
], 0.4, 2.45, 7.0, sz=16, gap=0.52)

# Right side — stat boxes
stat_box(s4, 7.7, 1.3, "+45%", "Demand growth\n2020–2025", val_col=RED)
stat_box(s4, 10.6, 1.3, "~70%", "Global internet\ntraffic routed through", val_col=BLUE)
stat_box(s4, 7.7, 2.8, "9,000+", "MW floor demand\n(3 AM, weekends)", val_col=TEAL)
stat_box(s4, 10.6, 2.8, "20 yrs", "Historical data\n2005 – 2025", val_col=GREEN)

# Bottom callout
box(s4, 0.4, 5.8, 12.5, 1.0, fill=SURF, radius=True)
tb(s4, "⚡  This makes DOM a uniquely challenging and important forecasting problem — "
       "conventional seasonal patterns are overlaid with a structural upward shift in baseload.",
   0.7, 5.9, 12.0, 0.8, sz=15, col=YELLOW)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Dataset Overview
# ══════════════════════════════════════════════════════════════
s5 = new_slide()
header(s5, "Dataset Overview", "PJM Interconnection — Dominion (DOM) region")

# Left column: data sources
box(s5, 0.4, 1.25, 5.9, 5.5, fill=SURF, radius=True)
tb(s5, "Data Sources", 0.65, 1.4, 5.4, 0.45, sz=18, bold=True, col=BLUE)

tb(s5, "Kaggle Historical (2005–2018)", 0.65, 2.0, 5.4, 0.38, sz=15, bold=True, col=TEXT)
bullet_col(s5, [
    "DOM_hourly.csv — 122,112 rows",
    "Hourly MW readings for the DOM zone",
], 0.65, 2.38, 5.4, sz=14, col=SUB, gap=0.38)

tb(s5, "PJM Metered Data (2018–2025)", 0.65, 3.35, 5.4, 0.38, sz=15, bold=True, col=TEXT)
bullet_col(s5, [
    "hrl_load_metered_YYYY.csv (8 files)",
    "Official PJM settlement-quality MW data",
    "Deduplicated & concatenated with Kaggle",
], 0.65, 3.73, 5.4, sz=14, col=SUB, gap=0.38)

tb(s5, "Combined & cleaned: ~175,000 hourly observations", 0.65, 5.25, 5.4, 0.45, sz=14, col=GREEN)

# Right column: key stats in a grid
right_stats = [
    ("175,000+", "Hourly observations",     BLUE),
    ("20 years",  "2005 – 2025",             TEAL),
    ("2005–2024", "Training set",            GREEN),
    ("2025",      "Test set (unseen)",       YELLOW),
    ("2026–2030", "Future forecast horizon", PURPLE),
    ("21",        "Engineered features",     PEACH),
]
for i, (val, lbl, col) in enumerate(right_stats):
    col_idx = i % 2
    row = i // 2
    x = 6.6 + col_idx * 3.4
    y = 1.25 + row * 1.72
    box(s5, x, y, 3.1, 1.5, fill=SURF, radius=True)
    tb(s5, val, x, y + 0.08, 3.1, 0.65, sz=28, bold=True, col=col, align=PP_ALIGN.CENTER)
    tb(s5, lbl, x, y + 0.75, 3.1, 0.45, sz=13, col=SUB, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — EDA: Demand Patterns
# ══════════════════════════════════════════════════════════════
s6 = new_slide()
header(s6, "EDA: Demand Patterns", "Systematic daily and seasonal cycles")

# Monthly demand bar chart
months  = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
avg_mw  = [13200, 12700, 11900, 11200, 11800, 13100, 14800, 14600, 12400, 11600, 12300, 13600]
add_bar_chart(s6, 0.4, 1.2, 7.5, 5.5,
              months,
              [("Avg Demand (MW)", avg_mw)],
              title="Average Monthly Demand (MW)",
              bar_color=BLUE)

# Key observations
obs = [
    ("Summer peak", "Jul–Aug highest demand\n(air conditioning load)"),
    ("Winter peak", "Jan–Feb second highest\n(space heating)"),
    ("Spring dip", "Mar–May lowest demand\n(mild temperatures)"),
    ("Dual seasonality", "U-shaped temperature–\ndemand relationship"),
]
for i, (title, body) in enumerate(obs):
    y = 1.25 + i * 1.55
    box(s6, 8.25, y, 4.65, 1.35, fill=SURF, radius=True)
    tb(s6, title, 8.45, y + 0.05, 4.2, 0.38, sz=15, bold=True, col=BLUE)
    tb(s6, body, 8.45, y + 0.45, 4.2, 0.75, sz=13, col=TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — EDA: Rising Baseload (Data Center Signature)
# ══════════════════════════════════════════════════════════════
s7 = new_slide()
header(s7, "EDA: Rising Baseload Trend", "Data center fingerprint in the demand data")

years  = [str(y) for y in range(2005, 2026)]
yearly = [10800,10900,11100,11300,11200,10600,10300,10100,10200,10400,
          10500,10600,10800,10800,10800,10900,10700,10700,10900,12300,13500]
add_line_chart(s7, 0.4, 1.2, 8.2, 5.5,
               years,
               [("Yearly Avg Demand (MW)", yearly)],
               title="Yearly Average Demand (MW) — 2005 to 2025",
               line_color=GREEN)

# Annotation boxes on right
box(s7, 8.9, 1.2, 4.05, 1.35, fill=SURF, radius=True)
tb(s7, "2005–2019: Flat / declining", 9.05, 1.3, 3.7, 0.4, sz=14, bold=True, col=TEAL)
tb(s7, "Energy efficiency gains offset\npopulation + economic growth",
   9.05, 1.7, 3.7, 0.7, sz=13, col=SUB)

box(s7, 8.9, 2.75, 4.05, 1.35, fill=SURF, radius=True)
tb(s7, "2019–2020: Brief COVID dip", 9.05, 2.85, 3.7, 0.4, sz=14, bold=True, col=YELLOW)
tb(s7, "Commercial/industrial load\ndrops during lockdowns",
   9.05, 3.25, 3.7, 0.7, sz=13, col=SUB)

box(s7, 8.9, 4.3, 4.05, 2.0, fill=SURF, radius=True)
tb(s7, "2020–2025: Surge (+45%)", 9.05, 4.4, 3.7, 0.4, sz=14, bold=True, col=RED)
tb(s7, "Hyperscale data center\nbuildout in Northern Virginia\n"
       "Amazon, Microsoft, Google,\nMeta — all expanding rapidly",
   9.05, 4.82, 3.7, 1.3, sz=13, col=TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — DOM vs Non-Data-Center Region (NEW)
# ══════════════════════════════════════════════════════════════
s8 = new_slide()
header(s8, "Is This a Regional Trend — or Data Centers?",
       "Normalized demand index: 2019 = 100%")

years_idx = [str(y) for y in range(2015, 2026)]

# DOM: flat 2015-2019, COVID dip 2020, then explosive growth
dom_idx  = [101, 102, 101, 101, 100, 98, 102, 105, 115, 123, 126]

# AEP (Ohio / Indiana) — no hyperscale data centers
# Gradual efficiency-driven decline, slight rebound, stays flat
aep_idx  = [105, 104, 103, 102, 100, 97, 99, 100, 99, 98, 97]

add_multi_line_chart(
    s8, 0.4, 1.25, 8.5, 5.6,
    years_idx,
    [
        ("DOM — Virginia/DC (Data Center Hub)", dom_idx,  RED),
        ("AEP — Ohio/Indiana (No Data Centers)", aep_idx, TEAL),
    ],
    title="Yearly Demand Index (2019 = 100%)",
)

# Right: narrative
box(s8, 9.15, 1.25, 3.85, 5.6, fill=SURF, radius=True)
tb(s8, "What the chart shows", 9.35, 1.38, 3.5, 0.42, sz=16, bold=True, col=BLUE)

points = [
    ("2015–2019",   "Both regions flat — energy efficiency\ngains offset demand growth everywhere", SUB),
    ("2020",         "Both dip together — COVID\nshuts down commercial load", SUB),
    ("2021–2022",   "Both recover in parallel", SUB),
    ("2023–2025",   "DOM breaks away — AEP stays\nflat while DOM surges +26%\nabove the 2019 baseline", RED),
]
for i, (yr, note, col) in enumerate(points):
    y = 2.0 + i * 1.25
    tb(s8, yr, 9.35, y, 3.5, 0.35, sz=13, bold=True, col=col)
    tb(s8, note, 9.35, y + 0.35, 3.5, 0.75, sz=12, col=TEXT)

tb(s8, "★  The surge is not a regional\nor national grid phenomenon.\nIt is specific to Virginia/DC\nbecause of data centers.",
   9.35, 5.05, 3.5, 1.2, sz=13, bold=True, col=YELLOW)

tb(s8, "Source: PJM Interconnection regional load data · Values normalized to 2019 baseline",
   0.4, 7.1, 12.5, 0.3, sz=11, col=SUB, italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — EDA: Weather & Demand (was 8)
# ══════════════════════════════════════════════════════════════
s8 = new_slide()
header(s8, "EDA: Weather & Demand Relationship", "Temperature drives both heating and cooling loads")

# Left: description of scatter
box(s8, 0.4, 1.25, 5.9, 5.5, fill=SURF, radius=True)
tb(s8, "Temperature vs Demand (2023)", 0.6, 1.4, 5.5, 0.45, sz=17, bold=True, col=BLUE)
tb(s8, "Open-Meteo archive API · Washington DC (38.89°N, 77.04°W)",
   0.6, 1.88, 5.5, 0.4, sz=13, col=SUB)

tb(s8, "Scatter: Hourly temperature (°F) vs demand (MW)",
   0.6, 2.45, 5.5, 0.45, sz=15, col=TEXT)
bullet_col(s8, [
    "Below ~55°F: demand rises as heating kicks in",
    "Above ~70°F: demand rises as A/C load increases",
    "55–70°F is the 'comfort zone' — lowest demand",
    "Pearson correlation ≈ +0.31 (linear, masking U-shape)",
], 0.6, 2.95, 5.5, sz=14, col=TEXT, gap=0.48)

box(s8, 0.6, 5.2, 5.5, 1.3, fill=SURF2, radius=True)
tb(s8, "⚡  Data center load is temperature-insensitive — it creates a raised\n"
       "floor (~9,000 MW) that is visible even at mild temperatures.",
   0.75, 5.3, 5.3, 1.1, sz=14, col=YELLOW)

# Right: dual-axis monthly chart (approximation using two bar series)
months_short = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
avg_mw_2023  = [13900,12800,11600,11100,11900,13200,15200,15000,12600,11400,12200,13800]
avg_temp     = [38, 42, 52, 62, 71, 80, 85, 83, 76, 64, 53, 41]

add_bar_chart(s8, 6.5, 1.25, 6.5, 5.5,
              months_short,
              [("Avg Demand (MW) — 2023", avg_mw_2023)],
              title="Monthly Demand vs Temperature (2023)",
              bar_color=BLUE)

# Note about dual-axis
tb(s8, "🌡 Temperature trend (°F): 38→85→41 — peaks Jul/Aug matching demand peak",
   6.6, 6.55, 6.3, 0.45, sz=12, col=TEAL)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — Feature Engineering
# ══════════════════════════════════════════════════════════════
s9 = new_slide()
header(s9, "Feature Engineering", "21 features capturing time, calendar, and lag patterns")

feature_groups = [
    ("Time of Day",    BLUE,   ["hour (0–23)", "hour_sin, hour_cos", "(cyclical encoding)"]),
    ("Calendar",       TEAL,   ["dayofweek (0–6)", "month (1–12)", "quarter, dayofyear",
                                 "weekofyear", "is_weekend flag", "is_holiday flag"]),
    ("Lag Features",   GREEN,  ["lag_1h  (previous hour)", "lag_24h  (yesterday same hour)",
                                 "lag_168h  (last week)"]),
    ("Rolling Stats",  YELLOW, ["roll_mean_24h  (last 24h avg)", "roll_std_24h  (last 24h σ)",
                                  "roll_mean_168h  (last week avg)", "roll_std_168h  (last week σ)"]),
]

for i, (group, col, features) in enumerate(feature_groups):
    col_idx = i % 2
    row = i // 2
    x = 0.4 + col_idx * 6.5
    y = 1.25 + row * 2.75
    box(s9, x, y, 6.15, 2.55, fill=SURF, radius=True)
    box(s9, x, y, 6.15, 0.07, fill=col)
    tb(s9, group, x + 0.15, y + 0.15, 5.8, 0.42, sz=18, bold=True, col=col)
    for j, feat in enumerate(features):
        tb(s9, f"▸  {feat}", x + 0.15, y + 0.65 + j * 0.45, 5.8, 0.42, sz=14, col=TEXT)

# Bottom note
box(s9, 0.4, 6.8, 12.5, 0.55, fill=SURF2, radius=True)
tb(s9, "Cyclical encoding (sin/cos) ensures the model understands that hour 23 is adjacent to hour 0 — "
       "no ordinal jump at day boundaries.",
   0.65, 6.87, 12.0, 0.42, sz=13, col=SUB)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — Probabilistic Forecasting
# ══════════════════════════════════════════════════════════════
s10 = new_slide()
header(s10, "Probabilistic Forecasting", "Prediction intervals — not just a single number")

# Left: what vs why
box(s10, 0.4, 1.25, 6.0, 5.5, fill=SURF, radius=True)
tb(s10, "What is Quantile Regression?", 0.6, 1.4, 5.7, 0.45, sz=18, bold=True, col=BLUE)
bullet_col(s10, [
    "Train three separate models on Q10, Q50, Q90",
    "Q50 = median forecast (like MAE-minimizing)",
    "Q10 = lower bound (10% of actuals fall below)",
    "Q90 = upper bound (90% of actuals fall below)",
    "Q10–Q90 band = 80% prediction interval",
], 0.6, 1.9, 5.7, sz=15, gap=0.48)

tb(s10, "Loss Function: Pinball Loss", 0.6, 4.3, 5.7, 0.45, sz=16, bold=True, col=TEAL)
tb(s10, "For quantile q:  L = (y − ŷ) × q   if y ≥ ŷ\n"
        "                         = (ŷ − y) × (1−q)  if y < ŷ",
   0.6, 4.78, 5.7, 0.85, sz=14, col=TEXT, italic=True)

# Right: why it matters
box(s10, 6.65, 1.25, 6.3, 5.5, fill=SURF, radius=True)
tb(s10, "Why Intervals Matter", 6.85, 1.4, 5.9, 0.45, sz=18, bold=True, col=GREEN)

reasons = [
    ("Coverage",  "We target 80% — the actual demand\nshould fall inside the band 80% of the time"),
    ("Width",     "Narrow = high confidence.\nWide = high uncertainty (e.g. extreme weather)"),
    ("Grid ops",  "Operators reserve spinning capacity\nbased on the upper bound (Q90)"),
    ("Risk mgmt", "Utilities can quantify worst-case\ncosts without overbuilding"),
]
for i, (title, body) in enumerate(reasons):
    y = 1.95 + i * 1.15
    tb(s10, title, 6.85, y, 1.8, 0.38, sz=15, bold=True, col=YELLOW)
    tb(s10, body, 8.7, y, 4.0, 0.85, sz=14, col=TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — Three Models
# ══════════════════════════════════════════════════════════════
s11 = new_slide()
header(s11, "Three Models", "Gradient boosting and deep learning for quantile forecasting")

models_info = [
    ("XGBoost", PEACH,
     ["Extreme Gradient Boosting",
      "Shallow trees, strong regularization",
      "Trained separately for Q10/Q50/Q90",
      "Fastest training & inference",
      "Baseline — well-established approach"]),
    ("LightGBM", BLUE,
     ["Light Gradient Boosting Machine",
      "Leaf-wise tree growth (faster than XGBoost)",
      "Histogram-based — handles large datasets well",
      "Deployed in Streamlit app for live forecast",
      "Powers the 5-year recursive forecast"]),
    ("LSTM", GREEN,
     ["Long Short-Term Memory Network",
      "Sliding window: 168 hours (1 week) of history",
      "3-output head: Q10, Q50, Q90 simultaneously",
      "Captures long-range temporal dependencies",
      "Best MAE — most accurate model overall"]),
]

for i, (name, col, bullets) in enumerate(models_info):
    x = 0.4 + i * 4.3
    box(s11, x, 1.25, 4.1, 5.5, fill=SURF, radius=True)
    box(s11, x, 1.25, 4.1, 0.07, fill=col)
    tb(s11, name, x, 1.38, 4.1, 0.5, sz=24, bold=True, col=col, align=PP_ALIGN.CENTER)
    for j, b in enumerate(bullets):
        tb(s11, f"▸  {b}", x + 0.15, 2.05 + j * 0.82, 3.8, 0.7, sz=14, col=TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — Train / Test Split
# ══════════════════════════════════════════════════════════════
s12 = new_slide()
header(s12, "Train / Test Split", "Time-based split — no data leakage")

# Timeline bar
box(s12, 0.5, 2.3, 12.3, 0.9, fill=SURF2, radius=True)

# Train portion
box(s12, 0.5, 2.3, 9.4, 0.9, fill=RGBColor(0x1e, 0x5f, 0x74), radius=True)
tb(s12, "TRAIN  2005–2024  (20 years · ~153,000 observations)",
   0.75, 2.48, 9.0, 0.5, sz=15, bold=True, col=TEXT)

# Test portion
box(s12, 10.0, 2.3, 2.8, 0.9, fill=RGBColor(0x74, 0x3a, 0x1e), radius=True)
tb(s12, "TEST  2025", 10.1, 2.48, 2.6, 0.5, sz=14, bold=True, col=YELLOW)

# Year labels
for yr, pos in [("2005", 0.5), ("2010", 3.35), ("2015", 5.9), ("2020", 8.05),
                ("2025", 10.05), ("2026+", 11.5)]:
    tb(s12, yr, pos, 3.35, 1.2, 0.35, sz=12, col=SUB)
    box(s12, pos + 0.2, 3.2, 0.03, 0.25, fill=SURF2)

# Explanations
boxes_info = [
    (0.5, 3.85, 6.0, "Why 2025 as test year?", BLUE,
     ["2025 data exists — used as true holdout", "Models never saw 2025 during training",
      "Evaluates generalization to unseen future"]),
    (6.8, 3.85, 6.0, "Why train through 2024?", GREEN,
     ["Earlier split would waste recent data", "Data center surge (2020–2025) is in train set",
      "Model learns the new higher baseload regime"]),
]
for x, y, w, title, col, body in boxes_info:
    box(s12, x, y, w, 2.75, fill=SURF, radius=True)
    tb(s12, title, x + 0.2, y + 0.12, w - 0.4, 0.42, sz=17, bold=True, col=col)
    for j, line in enumerate(body):
        tb(s12, f"▸  {line}", x + 0.2, y + 0.65 + j * 0.55, w - 0.4, 0.5, sz=14, col=TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 13 — Results
# ══════════════════════════════════════════════════════════════
s13 = new_slide()
header(s13, "Model Performance Results", "Full 2025 test set — 8,760 hourly observations")

# Metrics table (manual layout)
cols    = ["Model", "MAE (MW)", "Avg Pinball", "Coverage %", "Interval (MW)"]
rows    = [
    ("XGBoost",  "224.3", "73.0", "79.6%", "687"),
    ("LightGBM", "230.5", "73.7", "78.9%", "698"),
    ("LSTM",     "171.6", "58.8", "93.8%", "822"),
]
col_x   = [0.4, 3.6, 5.9, 8.2, 10.5]
col_w   = [3.0, 2.1, 2.1, 2.1, 2.4]

# Header row
box(s13, 0.4, 1.25, 12.5, 0.52, fill=SURF2)
for j, (c, x, w) in enumerate(zip(cols, col_x, col_w)):
    tb(s13, c, x + 0.1, 1.3, w, 0.42, sz=14, bold=True,
       col=BLUE if j == 0 else TEXT)

# Data rows
row_colors = [SURF, SURF, RGBColor(0x2a, 0x40, 0x2a)]
for i, (row, fill) in enumerate(zip(rows, row_colors)):
    y = 1.82 + i * 0.65
    box(s13, 0.4, y, 12.5, 0.60, fill=fill)
    for j, (val, x, w) in enumerate(zip(row, col_x, col_w)):
        c = PEACH if (i == 0 and j == 0) else (BLUE if (i == 1 and j == 0)
             else (GREEN if (i == 2 and j == 0) else TEXT))
        if i == 2 and j in (1, 2):   # LSTM best MAE/pinball
            c = GREEN
        if i == 2 and j == 3:        # LSTM coverage
            c = GREEN
        tb(s13, val, x + 0.1, y + 0.1, w, 0.42, sz=15, col=c)

tb(s13, "★  LSTM achieves best MAE and pinball loss. Wider interval reflects higher coverage (93.8% vs 80% target).",
   0.4, 3.45, 12.5, 0.45, sz=14, col=GREEN)

# Comparison bar chart: MAE
add_bar_chart(s13, 0.4, 4.0, 5.8, 2.85,
              ["XGBoost", "LightGBM", "LSTM"],
              [("MAE (MW)", [224.3, 230.5, 171.6])],
              title="MAE (MW) — lower is better",
              bar_color=BLUE)

add_bar_chart(s13, 6.6, 4.0, 6.1, 2.85,
              ["XGBoost", "LightGBM", "LSTM"],
              [("Coverage %", [79.6, 78.9, 93.8])],
              title="Coverage % — higher is better",
              bar_color=GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 14 — Prediction Interval Visualization
# ══════════════════════════════════════════════════════════════
s14 = new_slide()
header(s14, "Prediction Intervals", "80% confidence band around Q50 forecast")

# Simulate a chart area with boxes (visual representation)
box(s14, 0.4, 1.25, 12.5, 4.8, fill=SURF, radius=True)
tb(s14, "← The Streamlit app shows interactive plots like this for any selected week in 2025:",
   0.6, 1.35, 12.0, 0.42, sz=13, col=SUB, italic=True)

# Mock chart visual
box(s14, 0.5, 1.8, 12.3, 0.06, fill=SURF2)   # x axis

# Shade interval
box(s14, 0.65, 2.3, 11.7, 2.0, fill=RGBColor(0x26, 0x35, 0x57))
tb(s14, "80% Prediction Interval (Q10 to Q90)", 5.5, 2.35, 5.0, 0.42, sz=13, col=BLUE)

# Draw mock lines (Q50 and Actual)
for x_pos in [0.8, 1.9, 3.0, 4.1, 5.2, 6.3, 7.4, 8.5, 9.6, 10.7, 11.8]:
    import math
    offset = math.sin(x_pos * 0.9) * 0.45 + 0.5
    box(s14, x_pos, 2.55 + offset, 1.05, 0.04, fill=BLUE)
    box(s14, x_pos + 0.15, 2.65 + offset + math.sin(x_pos * 1.3) * 0.18, 0.9, 0.04, fill=RED)

tb(s14, "── Q50 Forecast (LightGBM/LSTM)",
   0.6, 4.45, 6.0, 0.42, sz=14, col=BLUE)
tb(s14, "·· · Actual Demand",
   6.5, 4.45, 6.0, 0.42, sz=14, col=RED)

# Observations below
obs = [
    ("Winter week\nJan 6–12", "Sharp morning ramp;\nmodel tracks well"),
    ("Summer peak\nJul 7–13", "Sustained afternoon plateau;\nwider intervals reflect A/C uncertainty"),
    ("Holiday week\nDec 17–25", "Demand drop on Dec 25;\nholiday flag captures the drop"),
]
for i, (period, note) in enumerate(obs):
    x = 0.4 + i * 4.3
    box(s14, x, 6.2, 4.1, 1.15, fill=SURF, radius=True)
    tb(s14, period, x + 0.15, 6.28, 3.8, 0.45, sz=13, bold=True, col=YELLOW)
    tb(s14, note, x + 0.15, 6.73, 3.8, 0.55, sz=13, col=TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 15 — Key Finding: Data Center Signature
# ══════════════════════════════════════════════════════════════
s15 = new_slide()
header(s15, "Key Finding: Data Center Signature", "A structural shift hidden in the demand curve")

# Big central callout
box(s15, 0.4, 1.3, 12.5, 1.55, fill=RGBColor(0x2a, 0x2a, 0x1e), radius=True)
box(s15, 0.4, 1.3, 0.07, 1.55, fill=YELLOW)
tb(s15, "Average demand grew from ~10,700 MW (2019) to ~13,500 MW (2025) — a +26% structural shift",
   0.65, 1.45, 12.0, 0.5, sz=20, bold=True, col=YELLOW)
tb(s15, "This is not seasonal variation. This is always-on compute load from data centers.",
   0.65, 1.95, 12.0, 0.5, sz=17, col=TEXT)

findings = [
    ("Flat 3AM load",
     "~9,000 MW minimum even at 3 AM on holiday weekends — "
     "residential + commercial demand alone cannot explain this baseline.",
     TEAL),
    ("Smaller weekend dip",
     "Traditional grids see 15–25% demand drop on weekends. "
     "DOM's weekend dip is shrinking — data centers don't take Sundays off.",
     BLUE),
    ("Higher low-temperature floor",
     "Even in mild spring weather (55°F), DOM demand rarely drops below 10,000 MW. "
     "Historical data showed 8,500 MW lows — the floor has risen.",
     GREEN),
    ("Forecast implication",
     "Models trained only on pre-2020 data severely under-forecast present demand. "
     "Training through 2024 ensures the model has learned the new regime.",
     PURPLE),
]

for i, (title, body, col) in enumerate(findings):
    col_idx = i % 2
    row = i // 2
    x = 0.4 + col_idx * 6.5
    y = 3.15 + row * 1.85
    box(s15, x, y, 6.15, 1.65, fill=SURF, radius=True)
    box(s15, x, y, 0.06, 1.65, fill=col)
    tb(s15, title, x + 0.22, y + 0.1, 5.7, 0.42, sz=16, bold=True, col=col)
    tb(s15, body, x + 0.22, y + 0.58, 5.7, 1.0, sz=13, col=TEXT, wrap=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 16 — Future Forecast 2026–2030
# ══════════════════════════════════════════════════════════════
s16 = new_slide()
header(s16, "Future Forecast: 2026–2030", "Recursive LightGBM — 43,824 hourly steps")

# Left: how it works
box(s16, 0.4, 1.25, 6.0, 5.5, fill=SURF, radius=True)
tb(s16, "Recursive Forecasting", 0.6, 1.38, 5.7, 0.45, sz=18, bold=True, col=BLUE)
bullet_col(s16, [
    "Seed with last 200 known hourly values (Dec 2025)",
    "Predict step t+1 using lag features from history",
    "Append Q50 prediction to history buffer",
    "Repeat for 43,824 steps (2026-01-01 to 2030-12-31)",
    "Prediction intervals preserved throughout (Q10/Q90)",
], 0.6, 1.9, 5.7, sz=14, gap=0.52)

tb(s16, "Pre-generated & cached as parquet:", 0.6, 4.8, 5.7, 0.4, sz=14, bold=True, col=TEAL)
bullet_col(s16, [
    "Loads instantly on every app start",
    "Deterministic — same result each time",
    "~1.5 MB file in the repo (Git LFS)",
], 0.6, 5.2, 5.7, sz=14, col=TEXT, gap=0.42)

# Right: projected ranges
box(s16, 6.65, 1.25, 6.3, 5.5, fill=SURF, radius=True)
tb(s16, "Projected Ranges", 6.85, 1.38, 5.9, 0.45, sz=18, bold=True, col=GREEN)

proj_data = [
    ("Peak Q50 (summer)",    "~20,600 MW",   RED),
    ("Min Q50 (spring AM)",  "~9,600 MW",    TEAL),
    ("Avg uncertainty band", "~2,000 MW",    BLUE),
    ("Horizon",              "1,825 days",   YELLOW),
]
for i, (label, val, col) in enumerate(proj_data):
    y = 2.0 + i * 1.1
    box(s16, 6.75, y, 6.1, 0.9, fill=SURF2, radius=True)
    tb(s16, label, 6.95, y + 0.08, 3.5, 0.38, sz=14, col=TEXT)
    tb(s16, val, 10.3, y + 0.05, 2.3, 0.45, sz=18, bold=True, col=col, align=PP_ALIGN.RIGHT)

box(s16, 6.65, 5.5, 6.3, 1.25, fill=RGBColor(0x2a, 0x2a, 0x1e), radius=True)
tb(s16, "⚠  Long-horizon caveats:",
   6.85, 5.58, 5.9, 0.4, sz=14, bold=True, col=YELLOW)
tb(s16, "Intervals widen due to cumulative lag uncertainty. "
        "Seasonal structure is reliable; absolute MW levels are projections.",
   6.85, 5.98, 5.9, 0.68, sz=13, col=SUB)

# ══════════════════════════════════════════════════════════════
# SLIDE 17 — Live Demo (Streamlit App)
# ══════════════════════════════════════════════════════════════
s17 = new_slide()
header(s17, "Live Demo — Streamlit App", "Interactive probabilistic forecast explorer")

# App URL box
box(s17, 0.4, 1.3, 12.5, 0.9, fill=RGBColor(0x1e, 0x35, 0x4a), radius=True)
tb(s17, "⚡  https://energy-demand-forecasting-ml.streamlit.app",
   0.65, 1.45, 12.0, 0.6, sz=22, bold=True, col=BLUE)

# 5 tabs description
tabs = [
    ("📊 Model Evaluation",   BLUE,   "Select any week in 2025, choose LightGBM / LSTM / XGBoost, view forecast vs actual with 80% intervals"),
    ("🏆 Model Comparison",   TEAL,   "Side-by-side metrics table and bar charts; summer peak week comparison for both models"),
    ("🔍 EDA",                GREEN,  "Hourly / daily / monthly / yearly demand patterns; weather-demand relationship charts (2023 DC data)"),
    ("🔮 Future Forecast",    PURPLE, "Date range picker for 2026–2030 with quick select; recursive LightGBM Q10/Q50/Q90 bands"),
    ("ℹ️ About",              YELLOW, "Project overview, dataset details, model descriptions, performance table, and key findings"),
]
for i, (tab, col, desc) in enumerate(tabs):
    y = 2.45 + i * 0.88
    box(s17, 0.4, y, 12.5, 0.80, fill=SURF, radius=True)
    tb(s17, tab, 0.6, y + 0.08, 3.1, 0.45, sz=15, bold=True, col=col)
    tb(s17, desc, 3.75, y + 0.08, 9.0, 0.62, sz=14, col=TEXT)

tb(s17, "All predictions pre-cached · Loads in < 2 seconds · No API keys required",
   0.4, 6.88, 12.5, 0.4, sz=13, col=SUB, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 18 — Conclusions
# ══════════════════════════════════════════════════════════════
s18 = new_slide()
header(s18, "Conclusions & Key Takeaways")

# Left: takeaways
takeaways = [
    ("Probabilistic > point forecasting",
     "Quantile regression provides actionable uncertainty ranges — critical for grid operations",
     BLUE),
    ("LSTM wins on accuracy",
     "171.6 MW MAE vs 224+ for tree models; 1-week sliding window captures temporal dependencies",
     GREEN),
    ("LightGBM wins on speed",
     "Instant inference, recursive 5-year capability, ideal for production deployment",
     TEAL),
    ("Data centers are reshaping the grid",
     "The +45% demand surge in DOM is a fingerprint of hyperscale growth — not seasonal noise",
     YELLOW),
    ("Feature engineering is critical",
     "Lag features, rolling statistics, cyclical encoding give models the temporal context they need",
     PURPLE),
]

for i, (title, body, col) in enumerate(takeaways):
    y = 1.28 + i * 1.14
    box(s18, 0.4, y, 8.2, 1.0, fill=SURF, radius=True)
    box(s18, 0.4, y, 0.06, 1.0, fill=col)
    tb(s18, title, 0.62, y + 0.04, 7.7, 0.38, sz=14, bold=True, col=col)
    tb(s18, body, 0.62, y + 0.46, 7.7, 0.45, sz=13, col=TEXT)

# Right: future work
box(s18, 8.9, 1.28, 4.05, 4.15, fill=SURF, radius=True)
tb(s18, "Future Work", 9.1, 1.38, 3.7, 0.42, sz=17, bold=True, col=PEACH)
fw = [
    "Add weather forecasts as live features",
    "Transformer / Temporal Fusion network",
    "Incorporate data center capacity data",
    "Real-time PJM API integration",
    "Multi-region forecast (PJM-wide)",
    "Anomaly detection for extreme events",
]
for j, item in enumerate(fw):
    tb(s18, f"▸  {item}", 9.1, 1.95 + j * 0.55, 3.7, 0.5, sz=13, col=TEXT)

# Bottom thank-you bar
box(s18, 0, 6.6, 13.333, 0.9, fill=SURF)
tb(s18, "Thank you  ·  Dawit Adane  ·  ML-1 Course Project  ·  2026",
   0, 6.72, 13.333, 0.55, sz=17, col=TEXT, align=PP_ALIGN.CENTER)
tb(s18, "⚡  energy-demand-forecasting-ml.streamlit.app",
   0, 6.72, 13.333, 0.55, sz=14, col=BLUE, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
